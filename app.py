import streamlit as st
import fitz
import pytesseract
import re
import time

from docx import Document
from pptx import Presentation
from pdf2image import convert_from_bytes

# ----------------------------------
# TESSERACT PATH
# ----------------------------------
pytesseract.pytesseract.tesseract_cmd = (
    r"C:\Program Files\Tesseract-OCR\tesseract.exe"
)

# ----------------------------------
# PAGE CONFIG
# ----------------------------------
st.set_page_config(
    page_title="Kindle Reading Assistant",
    layout="wide"
)

# Set up Session States for Immersion Engines
if "zen_mode" not in st.session_state:
    st.session_state.zen_mode = False
if "auto_scroll" not in st.session_state:
    st.session_state.auto_scroll = False

# Load true Variable-Weight Web Fonts across the full spectrum (300 to 900 weight)
st.markdown("""
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Lora:ital,wght@0,300..900;1,300..900&family=Merriweather:ital,wght@0,300..900;1,300..900&family=Playfair+Display:ital,wght@0,400..900;1,400..900&family=Courier+Prime:ital,wght@0,400;0,700;1,400;1,700&family=Bitter:ital,wght@0,300..900;1,300..900&display=swap');
        
        .block-container { padding-top: 1rem; padding-bottom: 1rem; max-width: 1000px !important; }
        div[data-testid="stSidebarCollapseButton"] { display: none; }
        
        /* Kindle publishing paragraph rules */
        .kindle-p {
            margin-top: 0px !important;
            margin-bottom: 0px !important;
            text-indent: 2em;
        }
        .kindle-p:first-of-type {
            text-indent: 0px !important;
        }
    </style>
""", unsafe_allow_html=True)

# ----------------------------------
# LAYOUT & CHUNKING ENGINE
# ----------------------------------
def clean_and_chunk_text(text, words_per_page=250):
    if not text.strip():
        return []
    
    text = text.replace("\r\n", "\n")
    paragraphs_raw = text.split("\n\n")
    cleaned_paragraphs = []
    
    for p in paragraphs_raw:
        cleaned_p = re.sub(r'(?<!\n)\n(?!\n)', ' ', p)
        cleaned_p = re.sub(r'\s+', ' ', cleaned_p).strip()
        if cleaned_p:
            cleaned_paragraphs.append(cleaned_p)

    full_clean_text = " ".join(cleaned_paragraphs)
    words = full_clean_text.split(" ")
    
    chunks = []
    for i in range(0, len(words), words_per_page):
        chunk = " ".join(words[i:i + words_per_page])
        if chunk.strip():
            chunk_formatted = re.sub(r'(?<=[.!?])\s+(?=[A-Z“])', '</p><p class="kindle-p">', chunk)
            chunks.append(f'<p class="kindle-p">{chunk_formatted}</p>')
            
    return chunks

# ----------------------------------
# PARSING FILE METHODS
# ----------------------------------
def ocr_pdf(pdf_file):
    text = ""
    try:
        pdf_file.seek(0)
        images = convert_from_bytes(pdf_file.read())
        for image in images:
            text += pytesseract.image_to_string(image) + "\n"
    except Exception as e:
        st.warning(f"OCR failed: {str(e)}")
    return text

def extract_pdf(file):
    text = ""
    try:
        file.seek(0)
        pdf = fitz.open(stream=file.read(), filetype="pdf")
        for page in pdf:
            text += page.get_text() + "\n\n"
        pdf.close()
        if not text.strip():
            text = ocr_pdf(file)
    except Exception as e:
        st.error(f"PDF Extraction Error: {str(e)}")
    return text

def extract_docx(file):
    try:
        doc = Document(file)
        return "\n\n".join(para.text for para in doc.paragraphs if para.text.strip())
    except Exception as e:
        return ""

def extract_pptx(file):
    text = ""
    try:
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n\n"
    except Exception as e:
        return ""

def extract_txt(file):
    try:
        return file.read().decode("utf-8", errors="ignore")
    except Exception as e:
        return ""

# ----------------------------------
# SIDEBAR DESIGN ENGINE (VARIABLE SPECTRUMS)
# ----------------------------------
if not st.session_state.zen_mode:
    st.sidebar.title("📖 Book Typography")
    font_display_name = st.sidebar.selectbox(
        "Typeface Design",
        ["Lora (Editorial Serif)", "Merriweather (Thick Book-serif)", "Bitter (Highly Legible Slab)", "Playfair (High-Contrast Display)", "Courier Prime (Monospace)"]
    )
    
    # Precise responsive layout limits
    font_size = st.sidebar.slider("Font Size (px)", 14, 50, 22)
    font_weight = st.sidebar.slider("Letter Thickness (Weight)", 300, 900, 400, step=50)
    line_height = st.sidebar.slider("Line Spacing", 1.2, 3.0, 1.8, step=0.1)
    letter_spacing = st.sidebar.slider("Tracking (Letter Gap)", -2, 8, 0)
    max_container_width = st.sidebar.slider("Page Width Margins", 450, 900, 650, step=50)

    st.sidebar.markdown("---")
    st.sidebar.subheader("🎨 Paper Stock Finish")
    theme = st.sidebar.selectbox("Palette Presets", ["Vintage Cream", "Textured Antique", "Bleached Pulp", "Charcoal Slate"])
    words_per_page = st.sidebar.slider("Words Per Page Density", 100, 500, 240, step=20)
else:
    # Stable fallback variables while hidden
    font_display_name = "Lora (Editorial Serif)"
    font_size, font_weight, line_height, letter_spacing, max_container_width = 22, 400, 1.8, 0, 650
    theme = "Vintage Cream"
    words_per_page = 240

# Theme Configuration Dictionary
theme_presets = {
    "Vintage Cream": {"bg": "#EFE6D5", "text": "#2B221E", "card_bg": "#FAF4E8", "grain": "rgba(43,34,30,0.015)"},
    "Textured Antique": {"bg": "#E6DCC3", "text": "#332A24", "card_bg": "#F5EFE0", "grain": "rgba(51,42,36,0.02)"},
    "Bleached Pulp": {"bg": "#EAECEE", "text": "#1A1A1A", "card_bg": "#FFFFFF", "grain": "rgba(0,0,0,0.01)"},
    "Charcoal Slate": {"bg": "#1A1B1C", "text": "#D5D8DC", "card_bg": "#252729", "grain": "rgba(255,255,255,0.015)"}
}
style_config = theme_presets[theme]

# Mapping strictly to loaded Google Variable Font stacks
font_map = {
    "Lora (Editorial Serif)": "'Lora', Georgia, serif",
    "Merriweather (Thick Book-serif)": "'Merriweather', serif",
    "Bitter (Highly Legible Slab)": "'Bitter', serif",
    "Playfair (High-Contrast Display)": "'Playfair Display', serif",
    "Courier Prime (Monospace)": "'Courier Prime', monospace"
}
font_family = font_map[font_display_name]

# ----------------------------------
# APP CORE EXECUTION
# ----------------------------------
if not st.session_state.zen_mode:
    uploaded_file = st.file_uploader("Drop document to parse...", type=["pdf", "docx", "pptx", "txt"], label_visibility="collapsed")
else:
    uploaded_file = st.session_state.get("uploaded_file", None)

if uploaded_file is not None:
    st.session_state["uploaded_file"] = uploaded_file
    state_config_hash = f"{uploaded_file.name}_{words_per_page}"

    if "state_hash" not in st.session_state or st.session_state.state_hash != state_config_hash:
        st.session_state.state_hash = state_config_hash
        st.session_state.current_para = 0
        
        ext = uploaded_file.name.split(".")[-1].lower()
        raw_text = ""
        if ext == "pdf": raw_text = extract_pdf(uploaded_file)
        elif ext == "docx": raw_text = extract_docx(uploaded_file)
        elif ext == "pptx": raw_text = extract_pptx(uploaded_file)
        elif ext == "txt": raw_text = extract_txt(uploaded_file)

        st.session_state.paragraphs = clean_and_chunk_text(raw_text, words_per_page)

    pages = st.session_state.paragraphs

    if not pages:
        st.error("No extractable content loaded.")
        st.stop()

    if st.session_state.current_para >= len(pages):
        st.session_state.current_para = 0

    current_idx = st.session_state.current_para
    current_page_text = pages[current_idx]

    total_remaining_pages = len(pages) - (current_idx + 1)
    minutes_left = max(1, int((total_remaining_pages * words_per_page) / 220))
    progress_percentage = int(((current_idx + 1) / len(pages)) * 100)

    # ----------------------------
    # CONTROL TOP-DECK BAR
    # ----------------------------
    ctrl_col1, ctrl_col2, ctrl_col3, ctrl_col4 = st.columns([2, 2, 2, 2])
    with ctrl_col1:
        if st.button("👁️ Toggle Focus Mode", use_container_width=True):
            st.session_state.zen_mode = not st.session_state.zen_mode
            st.rerun()
    with ctrl_col2:
        if st.button("⏱️ Start/Stop Auto-Page", use_container_width=True):
            st.session_state.auto_scroll = not st.session_state.auto_scroll
            st.rerun()
    with ctrl_col3:
        scrolling_pace = st.slider("Page Pace (Secs)", 10, 120, 45, disabled=not st.session_state.auto_scroll, label_visibility="collapsed")
    with ctrl_col4:
        soundscape = st.selectbox("🎶 Soundscape", ["Silent Focus", "Rainstorm Loop", "Cosmic Ambient"], label_visibility="collapsed")

    # Audio Engine Frame
    if soundscape != "Silent Focus":
        sound_urls = {
            "Rainstorm Loop": "https://assets.mixkit.co/active_storage/sfx/2433/2433-84.wav",
            "Cosmic Ambient": "https://www.soundhelix.com/examples/mp3/SoundHelix-Song-1.mp3"
        }
        st.components.v1.html(f"""
            <audio src="{sound_urls[soundscape]}" autoplay loop></audio>
            <script>document.querySelector('audio').volume = 0.35;</script>
        """, height=0)

    # ----------------------------
    # IMMERSIVE PAPER CONTAINER
    # ----------------------------
    st.html(f"""
<div style="background-color: {style_config['bg']}; padding: 35px 15px; border-radius: 16px; display: flex; justify-content: center; box-shadow: inset 0 0 40px rgba(0,0,0,0.08);">
    <div style="background-color: {style_config['card_bg']}; 
                background-image: linear-gradient(90deg, {style_config['grain']} 1px, transparent 1px), linear-gradient(0deg, {style_config['grain']} 1px, transparent 1px);
                background-size: 3px 3px;
                width: 100%; max-width: {max_container_width}px; min-height: 700px; 
                padding: 55px 60px 40px 60px; 
                box-shadow: 0px 12px 30px rgba(0,0,0,0.07), inset -15px 0 20px rgba(0,0,0,0.015), inset 15px 0 20px rgba(0,0,0,0.015); 
                border-radius: 4px; display: flex; flex-direction: column; justify-content: space-between; 
                border-left: 1px solid rgba(0,0,0,0.03); border-right: 1px solid rgba(0,0,0,0.03); position: relative;">
        
        <div style="position: absolute; left: 0; top: 0; bottom: 0; width: 6px; background: linear-gradient(to right, rgba(0,0,0,0.02), transparent); pointer-events: none;"></div>

        <div style="font-family: {font_family} !important; font-size: {font_size}px !important; font-weight: {font_weight} !important; color: {style_config['text']}; line-height: {line_height} !important; letter-spacing: {letter_spacing}px !important; text-align: justify; text-justify: inter-word; hyphens: auto; display: block; clear: both;">
            {current_page_text}
        </div>

        <div style="display: flex; justify-content: space-between; align-items: center; font-family: 'Lora', Georgia, serif; font-size: 13px; font-style: italic; color: {style_config['text']}; opacity: 0.55; margin-top: 50px; border-top: 1px dashed rgba(0,0,0,0.1); padding-top: 18px;">
            <div>{minutes_left} mins remaining</div>
            <div style="font-family: monospace; font-style: normal;">Page {current_idx + 1} of {len(pages)}</div>
            <div>{progress_percentage}% parsed</div>
        </div>
    </div>
</div>
""")

    # ----------------------------
    # NAVIGATION TRACK INTERFACE
    # ----------------------------
    st.write("")
    nav_col1, nav_col2, nav_col3 = st.columns([2, 3, 2])
    with nav_col1:
        if st.button("◁ Previous Page", use_container_width=True):
            if st.session_state.current_para > 0:
                st.session_state.current_para -= 1
                st.rerun()
    with nav_col2:
        selected_page = st.slider("Jump Track", 1, len(pages), current_idx + 1, label_visibility="collapsed")
        if selected_page != current_idx + 1:
            st.session_state.current_para = selected_page - 1
            st.rerun()
    with nav_col3:
        if st.button("Next Page ▷", use_container_width=True):
            if st.session_state.current_para < len(pages) - 1:
                st.session_state.current_para += 1
                st.rerun()

    # Auto-scrolling processing loop
    if st.session_state.auto_scroll and st.session_state.current_para < len(pages) - 1:
        time.sleep(scrolling_pace)
        st.session_state.current_para += 1
        st.rerun()

    # ----------------------------
    # AUDIO SPEECH UTTERANCE
    # ----------------------------
    if not st.session_state.zen_mode:
        with st.expander("🔊 Screen Reader / Audio Assistance"):
            raw_text_extracted = re.sub('<[^<]+?>', '', current_page_text)
            safe_text = raw_text_extracted.replace("`", "'").replace("\\", "\\\\").replace("\n", " ")
            
            tts_html = f"""
            <script>
            function readText() {{
                var message = new SpeechSynthesisUtterance(`{safe_text}`);
                message.rate = 1.0; message.pitch = 1.0;
                window.speechSynthesis.cancel();
                window.speechSynthesis.speak(message);
            }}
            function stopReading() {{ window.speechSynthesis.cancel(); }}
            </script>
            <div style="display: flex; gap: 10px;">
                <button onclick="readText()" style="background-color: #007BFF; color: white; border: none; padding: 8px 16px; border-radius: 4px; cursor: pointer;">Play Audio</button>
                <button onclick="stopReading()" style="background-color: #6C757D; color: white; border: none; padding: 8px 16px; border-radius: 4px; cursor: pointer;">Pause</button>
            </div>
            """
            st.components.v1.html(tts_html, height=50)
else:
    st.info("Upload a document to turn this dashboard into an immersive Kindle book reader.") 